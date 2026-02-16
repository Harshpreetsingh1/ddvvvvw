from pptx import Presentation

# Create a presentation object
ppt = Presentation()

# Add title slide
slide = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "The 12 Principles of Animation"
subtitle.text = "Exploring the Foundation of Animation"

# Add overview slide
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Overview of the Principles"
content.text = (
    "1. Squash and Stretch\n"
    "2. Anticipation\n"
    "3. Staging\n"
    "4. Straight Ahead and Pose to Pose\n"
    "5. Follow Through and Overlapping Action\n"
    "6. Slow In and Slow Out\n"
    "7. Arcs\n"
    "8. Secondary Action\n"
    "9. Timing\n"
    "10. Exaggeration\n"
    "11. Solid Drawing\n"
    "12. Appeal"
)

# List of principles and their explanations
principles = [
    ("Squash and Stretch", "Adds flexibility and life to characters and objects.", 
     "Defines the rigidity or flexibility of objects, showing weight and volume in motion."),
    ("Anticipation", "Prepares the audience for an action to come.", 
     "Without anticipation, the action feels abrupt and lacks clarity."),
    ("Staging", "Directs the audience's attention to the key story elements.", 
     "The presentation of an idea so that it is unmistakably clear."),
    ("Straight Ahead and Pose to Pose", "Spontaneity vs. planned animation.", 
     "Combines two contrasting methods of animating."),
    ("Follow Through and Overlapping Action", "Continuation of motion after action.", 
     "Objects in motion tend to keep moving past their primary action."),
    ("Slow In and Slow Out", "Shows acceleration and deceleration in motion.", 
     "Softens transitions by spacing frames for smoothness."),
    ("Arcs", "Depicts natural motion paths.", 
     "Most actions follow an arc, not a straight line."),
    ("Secondary Action", "Adds complementary motions.", 
     "Enhances main action to support the story."),
    ("Timing", "Illustrates rhythm and pacing.", 
     "Correct timing makes animations feel natural."),
    ("Exaggeration", "Amplifies emotions and actions.", 
     "Enhances believability and impact."),
    ("Solid Drawing", "Creates realistic form and volume.", 
     "Characters should feel like they exist in a real space."),
    ("Appeal", "Makes characters visually engaging.", 
     "A character should attract viewers with a unique and compelling design."),
]

# Add slides for each principle
for principle, simple_explanation, famous_explanation in principles:
    # Title slide for the principle
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    title.text = principle
    
    # Explanation slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = f"Explanation of {principle}"
    content.text = f"{simple_explanation}\n\n{famous_explanation}"

# Add conclusion slide
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "The 12 Principles are the foundation of engaging and lifelike animations. Mastering them brings characters and stories to life."

# Save the presentation 
ppt.save("12_Principles_of_Animation.pptx")

print("Presentation created successfully!")
